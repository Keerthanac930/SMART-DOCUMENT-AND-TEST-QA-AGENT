"""
Enhanced Document processing module for handling various file types
Supports images (OCR), audio (speech-to-text), Google Drive, and more
"""
import os
import io
import tempfile
import logging
from typing import List, Dict, Any, Optional, Union
from pathlib import Path
import hashlib
import mimetypes
import json
from datetime import datetime

# Document processing
import PyPDF2
import docx
import openpyxl
from pptx import Presentation
import pandas as pd

# Image processing and OCR
from PIL import Image
import pytesseract
import cv2
import numpy as np
import easyocr

# Audio processing
import speech_recognition as sr
from pydub import AudioSegment
import librosa
import soundfile as sf

# Google Drive integration
from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload

# Configuration
from config import (
    SUPPORTED_FILE_TYPES, IMAGE_MAX_SIZE_MB, IMAGE_OCR_LANGUAGES,
    AUDIO_MAX_DURATION_SECONDS, AUDIO_SAMPLE_RATE, GOOGLE_DRIVE_SCOPES,
    GOOGLE_DRIVE_FOLDER_ID, FEATURES
)

logger = logging.getLogger(__name__)

class EnhancedDocumentProcessor:
    """
    Enhanced document processor supporting images, audio, Google Drive, and more file types
    """
    
    def __init__(self):
        self.supported_extensions = set(SUPPORTED_FILE_TYPES)
        self.ocr_reader = None
        self.speech_recognizer = sr.Recognizer()
        self.google_drive_service = None
        
        # Initialize OCR reader if image processing is enabled
        if FEATURES.get('image_processing', False):
            try:
                self.ocr_reader = easyocr.Reader(['en'])  # Initialize with English
                logger.info("EasyOCR initialized successfully")
            except Exception as e:
                logger.warning(f"Failed to initialize EasyOCR: {e}")
                self.ocr_reader = None
        
        # Initialize Google Drive service if enabled
        if FEATURES.get('google_drive', False):
            self._init_google_drive_service()
    
    def _init_google_drive_service(self):
        """Initialize Google Drive service"""
        try:
            creds = None
            token_path = 'token.json'
            credentials_path = 'credentials.json'
            
            # Load existing credentials
            if os.path.exists(token_path):
                creds = Credentials.from_authorized_user_file(token_path, GOOGLE_DRIVE_SCOPES)
            
            # If no valid credentials, get new ones
            if not creds or not creds.valid:
                if creds and creds.expired and creds.refresh_token:
                    creds.refresh(Request())
                elif os.path.exists(credentials_path):
                    flow = InstalledAppFlow.from_client_secrets_file(
                        credentials_path, GOOGLE_DRIVE_SCOPES)
                    creds = flow.run_local_server(port=0)
                
                # Save credentials for next run
                if creds:
                    with open(token_path, 'w') as token:
                        token.write(creds.to_json())
            
            if creds:
                self.google_drive_service = build('drive', 'v3', credentials=creds)
                logger.info("Google Drive service initialized successfully")
            else:
                logger.warning("Google Drive service initialization failed - no credentials")
                
        except Exception as e:
            logger.error(f"Error initializing Google Drive service: {e}")
            self.google_drive_service = None
    
    def process_document(self, file_path: str, file_content: Optional[bytes] = None, 
                        file_name: Optional[str] = None) -> Dict[str, Any]:
        """
        Process a document and extract text content
        
        Args:
            file_path (str): Path to the document file or URL
            file_content (bytes, optional): File content if processing from memory
            file_name (str, optional): File name if different from path
            
        Returns:
            Dict[str, Any]: Document metadata and content
        """
        try:
            # Handle Google Drive files
            if file_path.startswith('drive://'):
                return self._process_google_drive_file(file_path)
            
            # Handle file content from memory
            if file_content is not None:
                return self._process_file_content(file_content, file_name or file_path)
            
            # Handle regular file path
            file_path = Path(file_path)
            
            if not file_path.exists():
                raise FileNotFoundError(f"File not found: {file_path}")
            
            if file_path.suffix.lower() not in self.supported_extensions:
                raise ValueError(f"Unsupported file type: {file_path.suffix}")
            
            # Generate file hash for unique identification
            file_hash = self._generate_file_hash(file_path)
            
            # Extract text based on file type
            text_content = self._extract_text_by_type(file_path)
            
            # Clean and preprocess text
            cleaned_text = self._clean_text(text_content)
            
            return {
                'file_path': str(file_path),
                'file_name': file_path.name,
                'file_size': file_path.stat().st_size,
                'file_hash': file_hash,
                'file_type': file_path.suffix.lower(),
                'content': cleaned_text,
                'word_count': len(cleaned_text.split()),
                'char_count': len(cleaned_text),
                'processed_at': datetime.now().isoformat()
            }
            
        except Exception as e:
            logger.error(f"Error processing document {file_path}: {str(e)}")
            raise
    
    def _process_google_drive_file(self, drive_path: str) -> Dict[str, Any]:
        """Process a file from Google Drive"""
        if not self.google_drive_service:
            raise ValueError("Google Drive service not initialized")
        
        try:
            # Extract file ID from drive:// URL
            file_id = drive_path.replace('drive://', '')
            
            # Get file metadata
            file_metadata = self.google_drive_service.files().get(
                fileId=file_id, fields='id,name,mimeType,size,createdTime,modifiedTime'
            ).execute()
            
            # Download file content
            request = self.google_drive_service.files().get_media(fileId=file_id)
            file_content = io.BytesIO()
            downloader = MediaIoBaseDownload(file_content, request)
            
            done = False
            while done is False:
                status, done = downloader.next_chunk()
            
            file_content.seek(0)
            
            # Process the downloaded content
            return self._process_file_content(
                file_content.getvalue(),
                file_metadata['name']
            )
            
        except Exception as e:
            logger.error(f"Error processing Google Drive file {drive_path}: {str(e)}")
            raise
    
    def _process_file_content(self, file_content: bytes, file_name: str) -> Dict[str, Any]:
        """Process file content from memory"""
        try:
            # Save to temporary file for processing
            with tempfile.NamedTemporaryFile(delete=False, suffix=Path(file_name).suffix) as tmp_file:
                tmp_file.write(file_content)
                tmp_file_path = tmp_file.name
            
            try:
                # Generate file hash
                file_hash = hashlib.sha256(file_content).hexdigest()
                
                # Extract text based on file extension
                file_path = Path(file_name)
                text_content = self._extract_text_by_type(Path(tmp_file_path), file_name)
                
                # Clean and preprocess text
                cleaned_text = self._clean_text(text_content)
                
                return {
                    'file_path': file_name,
                    'file_name': file_name,
                    'file_size': len(file_content),
                    'file_hash': file_hash,
                    'file_type': file_path.suffix.lower(),
                    'content': cleaned_text,
                    'word_count': len(cleaned_text.split()),
                    'char_count': len(cleaned_text),
                    'processed_at': datetime.now().isoformat(),
                    'source': 'memory'
                }
                
            finally:
                # Clean up temporary file
                os.unlink(tmp_file_path)
                
        except Exception as e:
            logger.error(f"Error processing file content for {file_name}: {str(e)}")
            raise
    
    def _extract_text_by_type(self, file_path: Path, original_name: Optional[str] = None) -> str:
        """Extract text based on file type"""
        file_ext = file_path.suffix.lower()
        file_name = original_name or file_path.name
        
        try:
            if file_ext == '.pdf':
                return self._extract_pdf_text(file_path)
            elif file_ext == '.txt':
                return self._extract_txt_text(file_path)
            elif file_ext in ['.docx', '.doc']:
                return self._extract_docx_text(file_path)
            elif file_ext in ['.xlsx', '.xls']:
                return self._extract_excel_text(file_path)
            elif file_ext in ['.pptx', '.ppt']:
                return self._extract_powerpoint_text(file_path)
            elif file_ext == '.csv':
                return self._extract_csv_text(file_path)
            elif file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp']:
                return self._extract_image_text(file_path)
            elif file_ext in ['.mp3', '.wav', '.m4a', '.aac', '.ogg', '.flac', '.wma']:
                return self._extract_audio_text(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_ext}")
                
        except Exception as e:
            logger.error(f"Error extracting text from {file_name}: {str(e)}")
            raise
    
    def _extract_pdf_text(self, file_path: Path) -> str:
        """Extract text from PDF file"""
        try:
            text_content = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text.strip():
                            text_content += f"\n--- Page {page_num + 1} ---\n"
                            text_content += page_text
                    except Exception as e:
                        logger.warning(f"Error extracting text from page {page_num + 1}: {str(e)}")
                        continue
                
            return text_content
        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {str(e)}")
            raise
    
    def _extract_txt_text(self, file_path: Path) -> str:
        """Extract text from TXT file with encoding detection"""
        try:
            import chardet
            
            # Detect encoding
            with open(file_path, 'rb') as file:
                raw_data = file.read()
                encoding_info = chardet.detect(raw_data)
                encoding = encoding_info['encoding'] or 'utf-8'
            
            with open(file_path, 'r', encoding=encoding) as file:
                return file.read()
                
        except Exception as e:
            logger.error(f"Error processing TXT file {file_path}: {str(e)}")
            raise
    
    def _extract_docx_text(self, file_path: Path) -> str:
        """Extract text from DOCX file"""
        try:
            doc = docx.Document(file_path)
            text_content = ""
            
            for paragraph in doc.paragraphs:
                text_content += paragraph.text + "\n"
            
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text_content += cell.text + " "
                    text_content += "\n"
            
            return text_content
        except Exception as e:
            logger.error(f"Error processing DOCX file {file_path}: {str(e)}")
            raise
    
    def _extract_excel_text(self, file_path: Path) -> str:
        """Extract text from Excel file"""
        try:
            workbook = openpyxl.load_workbook(file_path)
            text_content = ""
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_content += f"\n--- Sheet: {sheet_name} ---\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join([str(cell) for cell in row if cell is not None])
                    if row_text.strip():
                        text_content += row_text + "\n"
            
            return text_content
        except Exception as e:
            logger.error(f"Error processing Excel file {file_path}: {str(e)}")
            raise
    
    def _extract_powerpoint_text(self, file_path: Path) -> str:
        """Extract text from PowerPoint file"""
        try:
            presentation = Presentation(file_path)
            text_content = ""
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                text_content += f"\n--- Slide {slide_num} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content += shape.text + "\n"
            
            return text_content
        except Exception as e:
            logger.error(f"Error processing PowerPoint file {file_path}: {str(e)}")
            raise
    
    def _extract_csv_text(self, file_path: Path) -> str:
        """Extract text from CSV file"""
        try:
            df = pd.read_csv(file_path)
            return df.to_string(index=False)
        except Exception as e:
            logger.error(f"Error processing CSV file {file_path}: {str(e)}")
            raise
    
    def _extract_image_text(self, file_path: Path) -> str:
        """Extract text from image using OCR"""
        if not FEATURES.get('image_processing', False):
            return "Image processing is disabled"
        
        try:
            # Load and preprocess image
            image = Image.open(file_path)
            
            # Check file size
            file_size_mb = os.path.getsize(file_path) / (1024 * 1024)
            if file_size_mb > IMAGE_MAX_SIZE_MB:
                logger.warning(f"Image file {file_path} is too large ({file_size_mb:.2f}MB)")
                image = image.resize((image.width // 2, image.height // 2))
            
            # Preprocess image for better OCR
            if FEATURES.get('advanced_ocr', False):
                image = self._preprocess_image_for_ocr(image)
            
            # Extract text using EasyOCR if available, otherwise Tesseract
            if self.ocr_reader:
                text_content = self._extract_text_easyocr(image)
            else:
                text_content = self._extract_text_tesseract(image)
            
            return text_content or "No text found in image"
            
        except Exception as e:
            logger.error(f"Error processing image {file_path}: {str(e)}")
            return f"Error processing image: {str(e)}"
    
    def _preprocess_image_for_ocr(self, image: Image.Image) -> Image.Image:
        """Preprocess image for better OCR results"""
        try:
            # Convert to numpy array
            img_array = np.array(image)
            
            # Convert to grayscale
            if len(img_array.shape) == 3:
                gray = cv2.cvtColor(img_array, cv2.COLOR_RGB2GRAY)
            else:
                gray = img_array
            
            # Apply denoising
            denoised = cv2.fastNlMeansDenoising(gray)
            
            # Apply thresholding
            _, thresh = cv2.threshold(denoised, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
            
            # Convert back to PIL Image
            return Image.fromarray(thresh)
            
        except Exception as e:
            logger.warning(f"Error preprocessing image: {e}")
            return image
    
    def _extract_text_easyocr(self, image: Image.Image) -> str:
        """Extract text using EasyOCR"""
        try:
            # Convert PIL image to numpy array
            img_array = np.array(image)
            
            # Extract text
            results = self.ocr_reader.readtext(img_array)
            
            # Combine all text
            text_parts = []
            for (bbox, text, confidence) in results:
                if confidence > 0.5:  # Filter low confidence results
                    text_parts.append(text)
            
            return " ".join(text_parts)
            
        except Exception as e:
            logger.error(f"Error with EasyOCR: {e}")
            return ""
    
    def _extract_text_tesseract(self, image: Image.Image) -> str:
        """Extract text using Tesseract OCR"""
        try:
            # Configure Tesseract
            custom_config = r'--oem 3 --psm 6'
            
            # Extract text
            text = pytesseract.image_to_string(image, config=custom_config)
            
            return text
            
        except Exception as e:
            logger.error(f"Error with Tesseract OCR: {e}")
            return ""
    
    def _extract_audio_text(self, file_path: Path) -> str:
        """Extract text from audio file using speech recognition"""
        if not FEATURES.get('audio_processing', False):
            return "Audio processing is disabled"
        
        try:
            # Load and preprocess audio
            audio = AudioSegment.from_file(str(file_path))
            
            # Check duration
            duration_seconds = len(audio) / 1000
            if duration_seconds > AUDIO_MAX_DURATION_SECONDS:
                logger.warning(f"Audio file {file_path} is too long ({duration_seconds:.1f}s)")
                audio = audio[:AUDIO_MAX_DURATION_SECONDS * 1000]
            
            # Convert to WAV format for speech recognition
            with tempfile.NamedTemporaryFile(suffix='.wav', delete=False) as tmp_file:
                audio.export(tmp_file.name, format='wav')
                
                try:
                    # Perform speech recognition
                    with sr.AudioFile(tmp_file.name) as source:
                        audio_data = self.speech_recognizer.record(source)
                    
                    # Try multiple recognition services
                    text = self._recognize_speech_multiple_engines(audio_data)
                    
                    return text or "No speech detected in audio"
                    
                finally:
                    os.unlink(tmp_file.name)
                    
        except Exception as e:
            logger.error(f"Error processing audio {file_path}: {str(e)}")
            return f"Error processing audio: {str(e)}"
    
    def _recognize_speech_multiple_engines(self, audio_data) -> str:
        """Try multiple speech recognition engines"""
        engines = [
            ('google', 'Google Speech Recognition'),
            ('sphinx', 'Sphinx (offline)')
        ]
        
        for engine, name in engines:
            try:
                if engine == 'google':
                    text = self.speech_recognizer.recognize_google(audio_data)
                elif engine == 'sphinx':
                    text = self.speech_recognizer.recognize_sphinx(audio_data)
                
                if text and text.strip():
                    logger.info(f"Speech recognition successful using {name}")
                    return text.strip()
                    
            except sr.UnknownValueError:
                logger.warning(f"Could not understand audio with {name}")
                continue
            except sr.RequestError as e:
                logger.warning(f"Error with {name}: {e}")
                continue
        
        return ""
    
    def _clean_text(self, text: str) -> str:
        """Clean and preprocess extracted text"""
        if not text:
            return ""
        
        # Remove excessive whitespace
        lines = text.split('\n')
        cleaned_lines = []
        
        for line in lines:
            # Remove leading/trailing whitespace
            line = line.strip()
            
            # Skip empty lines
            if not line:
                continue
            
            # Skip lines that are just page/slide separators
            if line.startswith('--- ') and line.endswith(' ---'):
                continue
            
            cleaned_lines.append(line)
        
        # Join lines with single newlines
        cleaned_text = '\n'.join(cleaned_lines)
        
        # Remove excessive spaces
        import re
        cleaned_text = re.sub(r'\s+', ' ', cleaned_text)
        
        return cleaned_text.strip()
    
    def _generate_file_hash(self, file_path: Path) -> str:
        """Generate SHA-256 hash of file for unique identification"""
        hash_sha256 = hashlib.sha256()
        
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_sha256.update(chunk)
        
        return hash_sha256.hexdigest()
    
    def chunk_text(self, text: str, chunk_size: int = 1000, overlap: int = 200) -> List[Dict[str, Any]]:
        """
        Split text into overlapping chunks for better processing
        
        Args:
            text (str): Text to chunk
            chunk_size (int): Size of each chunk
            overlap (int): Overlap between chunks
            
        Returns:
            List[Dict[str, Any]]: List of text chunks with metadata
        """
        if not text:
            return []
        
        chunks = []
        start = 0
        chunk_id = 0
        
        while start < len(text):
            end = start + chunk_size
            
            # Try to break at sentence boundary
            if end < len(text):
                # Look for sentence endings near the chunk boundary
                for i in range(end, max(start + chunk_size - 100, start), -1):
                    if text[i] in '.!?':
                        end = i + 1
                        break
            
            chunk_text = text[start:end].strip()
            
            if chunk_text:  # Only add non-empty chunks
                chunks.append({
                    'chunk_id': chunk_id,
                    'text': chunk_text,
                    'start_pos': start,
                    'end_pos': end,
                    'length': len(chunk_text)
                })
                chunk_id += 1
            
            # Move start position with overlap
            start = end - overlap
            if start <= 0:
                start = end
        
        return chunks
    
    def get_document_summary(self, document_data: Dict[str, Any]) -> Dict[str, Any]:
        """
        Generate a summary of the document
        
        Args:
            document_data (Dict[str, Any]): Document data from process_document
            
        Returns:
            Dict[str, Any]: Document summary
        """
        content = document_data['content']
        
        # Basic statistics
        words = content.split()
        sentences = content.split('.')
        paragraphs = content.split('\n\n')
        
        # Get first few sentences as preview
        preview = '. '.join(sentences[:3]) + '.' if len(sentences) > 3 else content[:200] + "..." if len(content) > 200 else content
        
        return {
            'file_name': document_data['file_name'],
            'file_type': document_data['file_type'],
            'file_size_mb': round(document_data['file_size'] / (1024 * 1024), 2),
            'word_count': document_data['word_count'],
            'char_count': document_data['char_count'],
            'sentence_count': len([s for s in sentences if s.strip()]),
            'paragraph_count': len([p for p in paragraphs if p.strip()]),
            'preview': preview,
            'processed_at': document_data.get('processed_at', 'Unknown'),
            'source': document_data.get('source', 'file')
        }
    
    def list_google_drive_files(self, folder_id: Optional[str] = None) -> List[Dict[str, Any]]:
        """List files from Google Drive folder"""
        if not self.google_drive_service:
            raise ValueError("Google Drive service not initialized")
        
        try:
            folder_id = folder_id or GOOGLE_DRIVE_FOLDER_ID
            
            if folder_id:
                query = f"'{folder_id}' in parents"
            else:
                query = "mimeType != 'application/vnd.google-apps.folder'"
            
            results = self.google_drive_service.files().list(
                q=query,
                fields="files(id,name,mimeType,size,createdTime,modifiedTime)",
                pageSize=100
            ).execute()
            
            files = []
            for file_info in results.get('files', []):
                # Only include supported file types
                file_ext = Path(file_info['name']).suffix.lower()
                if file_ext in self.supported_extensions:
                    files.append({
                        'id': file_info['id'],
                        'name': file_info['name'],
                        'mime_type': file_info.get('mimeType', ''),
                        'size': int(file_info.get('size', 0)),
                        'created_time': file_info.get('createdTime', ''),
                        'modified_time': file_info.get('modifiedTime', ''),
                        'drive_url': f"drive://{file_info['id']}"
                    })
            
            return files
            
        except Exception as e:
            logger.error(f"Error listing Google Drive files: {e}")
            raise
    
    def get_supported_file_types(self) -> List[str]:
        """Get list of supported file types"""
        return list(self.supported_extensions)
    
    def get_processing_capabilities(self) -> Dict[str, Any]:
        """Get information about processing capabilities"""
        return {
            'supported_file_types': self.get_supported_file_types(),
            'image_processing': FEATURES.get('image_processing', False),
            'audio_processing': FEATURES.get('audio_processing', False),
            'google_drive_enabled': self.google_drive_service is not None,
            'ocr_available': self.ocr_reader is not None,
            'speech_recognition_available': True,
            'max_file_size_mb': IMAGE_MAX_SIZE_MB,
            'max_audio_duration_seconds': AUDIO_MAX_DURATION_SECONDS
        }
